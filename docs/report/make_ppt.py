"""Build the CampusEase final-defense presentation.

Reproducible: re-run to regenerate campusease_presentation.pptx. Uses the
institutional template theme for background/fonts, but lays out every slide with
explicit, controlled geometry (manual title band + body box + fitted images) so
nothing clips or overlaps. Uses every generated diagram plus the application
screenshots. Requires python-pptx + Pillow.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
TEMPLATE = Path("/Users/nirajkafle/Desktop/niraj/dev-projects/pocket-drs/dump/BIT Project PPT Sample.pptx")
FIG = BASE / "figures"
OUT = BASE / "campusease_presentation.pptx"

NAVY = RGBColor(0x12, 0x33, 0x5B)
INK = RGBColor(0x20, 0x20, 0x20)
ACCENT = RGBColor(0x2E, 0x6F, 0xB5)

prs = Presentation(str(TEMPLATE))
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[10]

MARGIN = Inches(0.55)
TITLE_TOP = Inches(0.35)
TITLE_H = Inches(0.85)
BODY_TOP = Inches(1.35)
BODY_H = SH - BODY_TOP - Inches(0.95)  # leave room for the template footer logos
CONTENT_W = SW - 2 * MARGIN


def _new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return s


def _add_title(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(MARGIN, TITLE_TOP, CONTENT_W, TITLE_H)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    line = slide.shapes.add_shape(1, MARGIN, TITLE_TOP + TITLE_H, Inches(2.2), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background(); line.shadow.inherit = False


def bullet_slide(title: str, bullets: list[tuple[int, str]]) -> None:
    s = _new_slide()
    _add_title(s, title)
    tb = s.shapes.add_textbox(MARGIN, BODY_TOP, CONTENT_W, BODY_H)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, (level, text) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if level == 0 else "◦  ") + text
        p.level = level
        p.font.size = Pt(17 if level == 0 else 14)
        p.font.color.rgb = INK
        p.space_after = Pt(7 if level == 0 else 3)
        p.line_spacing = 1.05


def _fit(img: Path, max_w: int, max_h: int) -> tuple[int, int]:
    with Image.open(img) as im:
        w, h = im.size
    ar = w / h
    width, height = max_w, int(max_w / ar)
    if height > max_h:
        height, width = max_h, int(max_h * ar)
    return width, height


def image_slide(title: str, images: list[Path], captions: list[str] | None = None) -> None:
    s = _new_slide()
    _add_title(s, title)
    n = len(images)
    gap = Inches(0.3)
    cap_h = Inches(0.3) if captions else Inches(0.0)
    cell_w = int((CONTENT_W - gap * (n - 1)) / n)
    area_h = int(BODY_H - cap_h)
    x = MARGIN
    for k, img in enumerate(images):
        w, h = _fit(img, cell_w, area_h)
        s.shapes.add_picture(str(img), x + (cell_w - w) // 2, BODY_TOP + (area_h - h) // 2, width=w, height=h)
        if captions:
            cb = s.shapes.add_textbox(x, BODY_TOP + area_h, cell_w, cap_h)
            cp = cb.text_frame.paragraphs[0]
            cp.text = captions[k]; cp.alignment = PP_ALIGN.CENTER
            cp.font.size = Pt(12); cp.font.color.rgb = INK
        x += cell_w + gap


def title_slide() -> None:
    s = _new_slide()
    tf = s.shapes.add_textbox(MARGIN, Inches(1.3), CONTENT_W, Inches(3.0)).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rows = [
        ("CampusEase", 40, True, NAVY, 6),
        ("A Web-Based College Automation System", 18, False, INK, 0),
        ("for Academic and Administrative Management", 18, False, INK, 14),
        ("Bhupendra Malla, BIT (LC0003001648)", 16, False, INK, 2),
        ("Supervisor: Mr. Saishab Bhattarai", 15, False, INK, 2),
        ("Phoenix College of Management", 15, False, ACCENT, 0),
    ]
    for i, (text, size, bold, color, after) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
        p.space_after = Pt(after)


_sldIdLst = prs.slides._sldIdLst
for sid in list(_sldIdLst):
    prs.part.drop_rel(sid.get(qn("r:id")))
    _sldIdLst.remove(sid)

title_slide()

bullet_slide("Introduction", [
    (0, "Colleges run on many overlapping manual processes: enrollment, attendance, assignments, exams, fees, scheduling, and communication."),
    (0, "Handled on paper or across disconnected tools, this is slow, error-prone, and hard to track."),
    (1, "Records get duplicated or lost; attendance and marks are tedious to compile; updates are easy to miss."),
    (0, "CampusEase brings these activities onto a single web platform with role-based access for students, faculty, secretaries, and administrators."),
    (0, "Adds modern conveniences: face-recognition attendance, online fee payment, and real-time messaging."),
])

bullet_slide("Problem Statement", [
    (0, "Fragmented information: records spread across registers, spreadsheets, and unrelated tools, with no single source of truth."),
    (0, "Manual, repetitive work: attendance, internal marks, fee reconciliation, and reports done by hand."),
    (0, "Errors and inconsistency from re-entering the same data in many places, with no audit trail."),
    (0, "Weak communication: updates reach students through notice boards and informal channels."),
    (0, "Limited access control: hard to ensure each role sees and does only what it should."),
])

bullet_slide("Objectives", [
    (0, "General: a web-based college automation system centralising academic and administrative operations with role-based services."),
    (0, "Specific objectives:"),
    (1, "Design a role-based access model for students, faculty, secretaries, and administrators."),
    (1, "Implement enrollment, scheduling, and academic-record management."),
    (1, "Provide manual, OTP, and face-recognition attendance."),
    (1, "Automate assignments, marks entry, and internal-marks calculation."),
    (1, "Integrate online fee payment (Khalti) and real-time communication."),
])

bullet_slide("Functional Requirements", [
    (0, "Registration with email verification, authentication, and role-based access."),
    (0, "Subject enrollment (enrollment keys), class scheduling, and academic records."),
    (0, "Attendance by manual entry, one-time password, and face recognition."),
    (0, "Assignment posting and submission, marks entry, and internal-marks calculation."),
    (0, "Online fee payment via Khalti with administrative approval."),
    (0, "Real-time messaging, discussions, events, clubs, ID cards, and reports."),
])

bullet_slide("Non-functional Requirements", [
    (0, "Performance: respond within a couple of seconds; deliver messages in real time."),
    (0, "Usability: intuitive, role-specific interfaces with no special training."),
    (0, "Security: bcrypt-hashed passwords, JWT-authenticated requests, role-based access."),
    (0, "Reliability: handle invalid input and unavailable external services gracefully."),
    (0, "Maintainability and portability: modular code; any Node.js + MongoDB host."),
])

bullet_slide("Development Methodology", [
    (0, "Agile, sprint-based development suited to a system of many distinct modules."),
    (1, "Each sprint delivered working modules: data model, API, Angular screens, then testing."),
    (0, "Role model and shared user data model were built first, since later features depend on them."),
    (0, "Git and GitHub used throughout for incremental tracking and continuous integration."),
])

image_slide("Use-Case Model", [FIG / "use_case_diagram.png"])
image_slide("Data Modeling: ER Diagram", [FIG / "er_diagram.png"])
image_slide("Process Modeling: Data Flow Diagrams",
            [FIG / "dfd_level0.png", FIG / "dfd_level1.png"],
            ["Level 0 (context)", "Level 1"])
image_slide("System Architecture", [FIG / "architecture.png"])
image_slide("Key Process Models",
            [FIG / "activity_diagram.png", FIG / "sequence_diagram.png"],
            ["Activity: face attendance", "Sequence: online fee payment"])
image_slide("System Flowchart", [FIG / "flowchart.png"])

bullet_slide("Implementing Tools", [
    (0, "Angular + TypeScript + Bootstrap: single-page, role-specific frontend portals."),
    (0, "Node.js + Express: REST API; Socket.io for real-time messaging and presence."),
    (0, "MongoDB + Mongoose: NoSQL document database and object-document mapper."),
    (0, "JWT + bcrypt: stateless authentication and one-way password hashing."),
    (0, "Khalti (payment), Nodemailer (email/OTP), Multer (uploads), Highcharts (dashboards)."),
])

bullet_slide("Implementation: Module Details", [
    (0, "Authentication & Role: registration, email verify, login, JWT, role/department checks."),
    (0, "Academic: enrollment by key, class scheduling, academic records."),
    (0, "Attendance: manual, OTP, and face-recognition (webcam capture and matching)."),
    (0, "Assignments & Marks: posting/submission, marks entry, automatic internal-marks calc."),
    (0, "Fee Payment: Khalti online payment with admin approval workflow."),
    (0, "Communication & Admin: real-time chat, discussions, events, clubs, user management."),
])

image_slide("Application Demo",
            [FIG / "app_admin_dashboard.png", FIG / "app_user_management.png"],
            ["Admin console", "User management"])

bullet_slide("Testing", [
    (0, "Tested module by module, then end to end through the UI for each role."),
    (0, "Unit tests: auth/JWT/role checks, enrollment, attendance, internal-marks, fee, messaging."),
    (0, "System tests: register-verify-login, attendance by face/OTP, assignment flow, fee pay+approve, real-time chat, user management."),
    (0, "All unit and system test cases passed; role-based access reliably allowed/denied actions."),
])

bullet_slide("Conclusion", [
    (0, "A complete, working platform unifying academics, administration, and communication for a college."),
    (0, "Role-based access, three attendance methods, automatic internal marks, online fee payment, and real-time messaging."),
    (0, "Built on Angular + Node/Express + Socket.io + MongoDB, secured with JWT and bcrypt; all test cases passed."),
    (0, "Future work: mobile app, push notifications, advanced analytics, scalable deployment, more payment options."),
])

prs.save(str(OUT))
print(f"Wrote {OUT}  ({len(prs.slides)} slides)")
